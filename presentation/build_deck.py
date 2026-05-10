#!/usr/bin/env python3
"""Fill content into the Intro-to-CV pptx template.

The template ships with 5 slides:
    1. Title  (Yiqiao Liu / Taijia Liang)
    2. Related Work & Methodology      (section divider, empty body)
    3. Preliminary Results              (section divider, empty body)
    4. Analysis, Limitations & Next Steps  (section divider, empty body)
    5. Thank you

We insert 7 new content slides + fill the Analysis slide. Final 12-slide order:

    1.  Title                                       (Yiqiao + Taijia)
    2.  Related Work & Methodology  (divider)
    3.  OpenPose                                    (Yiqiao)
    4.  HMR                                         (Yiqiao)
    5.  HRNet                                       (Taijia)
    6.  From 2D to 3D — SMPL + reprojection         (Taijia)
    7.  From 3D to 4D — HMR 2.0 + PHALP             (Taijia)
    8.  Preliminary Results          (divider)
    9.  Image pipeline — 4 quad-plots               (Taijia)
    10. Video pipeline — 4D tennis                  (Taijia)
    11. Analysis, Limitations & Next Steps          (Taijia)
    12. Thank you

Design system:
    Primary head/body font:  Inter Display / Inter   (Helvetica Neue → Arial fallback)
    Mono accent (stats):     JetBrains Mono          (Consolas → Menlo fallback)
    Palette:  #330662 deep purple  ·  #C03960 coral  ·  #058DC7 cyan  ·  #FAF8FD wash
    Visual language:  numbered concept cards, hero stats, pill tags, accent bars.
"""
from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "presentation" / "template.pptx"
OUT = ROOT / "presentation" / "Intro to CV.pptx"
RESULTS = ROOT / "demo" / "results"

# ────────────────────────── design tokens ──────────────────────────

HEAD_FONT = "Inter Display"
BODY_FONT = "Inter"
MONO_FONT = "JetBrains Mono"

PURPLE = RGBColor(0x33, 0x06, 0x62)
PURPLE_SOFT = RGBColor(0x6E, 0x4C, 0x97)
CORAL = RGBColor(0xC0, 0x39, 0x60)
CYAN = RGBColor(0x05, 0x8D, 0xC7)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x88, 0x88, 0x88)
GREY_SOFT = RGBColor(0xB8, 0xB8, 0xB8)
WASH = RGBColor(0xFA, 0xF8, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ────────────────────────── primitives ──────────────────────────

def set_run(run, *, font=BODY_FONT, size=14, color=DARK, bold=False, italic=False, spacing=None):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    if spacing is not None:
        # XML manipulation: a:rPr/spc is character spacing in 100ths of a point
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(spacing))


def add_text(slide, left, top, width, height, paragraphs, *, anchor="t"):
    """paragraphs: list of (runs, paragraph_kwargs) where runs is list of (text, run_kwargs).

    Returns the textbox shape.
    """
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if anchor == "c":
        tf.vertical_anchor = 3  # MSO_ANCHOR.MIDDLE
    elif anchor == "b":
        tf.vertical_anchor = 4  # BOTTOM
    for i, (runs, pkw) in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pkw.get("align", PP_ALIGN.LEFT)
        if "space_before" in pkw:
            p.space_before = Pt(pkw["space_before"])
        if "space_after" in pkw:
            p.space_after = Pt(pkw["space_after"])
        if "line_spacing" in pkw:
            p.line_spacing = pkw["line_spacing"]
        for txt, rkw in runs:
            r = p.add_run()
            r.text = txt
            set_run(r, **rkw)
    return tb


def add_simple_text(slide, left, top, width, height, text, **kw):
    return add_text(slide, left, top, width, height,
                    [([(text, kw)], {"align": kw.pop("align", PP_ALIGN.LEFT)})])


def add_image(slide, path, left, top, width=None, height=None):
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                    Inches(width) if width else None,
                                    Inches(height) if height else None)


def add_line(slide, x1, y1, x2, y2, *, color=GREY_SOFT, weight=0.75):
    from pptx.enum.shapes import MSO_CONNECTOR
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, *, fill, line=None, line_w=0,
                    rounded=False, corner_radius=None):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    rect = slide.shapes.add_shape(shp, Inches(left), Inches(top), Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_w)
    # Suppress shadow on shapes (cleaner look)
    spPr = rect.fill._xPr
    effectLst = spPr.find(qn("a:effectLst"))
    if effectLst is None:
        effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    # Empty effectLst removes inherited shadow
    return rect


# ────────────────────────── components ──────────────────────────

def add_title(slide, text, *, kicker=None, color=PURPLE):
    """Hero title. Optional kicker (small label above) for visual rhythm."""
    if kicker:
        add_simple_text(slide, 0.36, 0.42, 9.2, 0.30, kicker.upper(),
                        font=BODY_FONT, size=10, color=GREY, bold=True, spacing=200)
        title_top = 0.74
    else:
        title_top = 0.50
    add_simple_text(slide, 0.36, title_top, 9.2, 0.70, text,
                    font=HEAD_FONT, size=26, color=color, bold=True, spacing=-15)


def add_eyebrow(slide, left, top, width, text, *, color=GREY, size=10):
    """A tiny label above content (e.g., 'CITATION', 'PARADIGM')."""
    add_simple_text(slide, left, top, width, 0.24, text.upper(),
                    font=BODY_FONT, size=size, color=color, bold=True, spacing=200)


def add_speaker_chip(slide, who):
    """Compact pill in the bottom-left."""
    L, T, W, H = 0.36, 5.20, 0.78, 0.26
    rect = add_filled_rect(slide, L, T, W, H, fill=PURPLE, rounded=True)
    rect.text_frame.margin_left = rect.text_frame.margin_right = Emu(0)
    rect.text_frame.margin_top = rect.text_frame.margin_bottom = Emu(0)
    p = rect.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = who
    set_run(r, font=BODY_FONT, size=10, color=WHITE, bold=True, spacing=100)


def add_pill(slide, left, top, width, height, text, *, fill=PURPLE, fg=WHITE, size=10, bold=True):
    rect = add_filled_rect(slide, left, top, width, height, fill=fill, rounded=True)
    rect.text_frame.margin_left = Inches(0.10)
    rect.text_frame.margin_right = Inches(0.10)
    rect.text_frame.margin_top = rect.text_frame.margin_bottom = Emu(0)
    p = rect.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    set_run(r, font=BODY_FONT, size=size, color=fg, bold=bold, spacing=80)
    return rect


def add_concept_card(slide, left, top, width, height, *, num, header, body, color=PURPLE):
    """Numbered concept card: tiny step number + bold header + small body text."""
    # Left accent bar
    add_filled_rect(slide, left, top, 0.06, height, fill=color)

    # Number
    add_simple_text(slide, left + 0.18, top + 0.04, 0.5, 0.32,
                    f"{num:02d}", font=HEAD_FONT, size=15, color=color, bold=True)
    # Header
    add_simple_text(slide, left + 0.74, top + 0.04, width - 0.85, 0.32,
                    header, font=HEAD_FONT, size=14, color=DARK, bold=True)
    # Body
    add_simple_text(slide, left + 0.74, top + 0.40, width - 0.85, height - 0.40,
                    body, font=BODY_FONT, size=11, color=GREY, italic=False)


def add_hero_stat(slide, left, top, width, *, big_text, label,
                  big_color=PURPLE, label_color=GREY, big_size=44, label_size=11):
    add_simple_text(slide, left, top, width, big_size * 0.018 + 0.2,
                    big_text, font=HEAD_FONT, size=big_size, color=big_color,
                    bold=True, spacing=-30, align=PP_ALIGN.CENTER)
    add_simple_text(slide, left, top + big_size * 0.018 + 0.1, width, 0.30,
                    label.upper(), font=BODY_FONT, size=label_size,
                    color=label_color, bold=True, spacing=200, align=PP_ALIGN.CENTER)


def add_divider(slide, left, top, width, *, color=GREY_SOFT, weight=0.5):
    add_line(slide, left, top, left + width, top, color=color, weight=weight)


def add_caption(slide, left, top, width, text, *, color=GREY, size=9, italic=True):
    add_simple_text(slide, left, top, width, 0.26, text,
                    font=BODY_FONT, size=size, color=color, italic=italic,
                    align=PP_ALIGN.CENTER)


def add_arrow(slide, x1, y1, x2, y2, *, color=PURPLE, weight=1.25):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                       Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    # Add arrowhead
    ln = conn.line._get_or_add_ln()
    head = etree.SubElement(ln, qn("a:headEnd"))
    head.set("type", "none")
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle"); tail.set("w", "med"); tail.set("len", "med")
    return conn


# ────────────────────────── slide builders ──────────────────────────

def build_openpose(slide):
    add_title(slide, "OpenPose", kicker="01  ·  2D Pose Estimation  ·  Bottom-up")
    add_speaker_chip(slide, "Yiqiao")

    # Citation in a pill on the right of the kicker line
    add_pill(slide, 7.20, 0.42, 2.40, 0.30, "Cao et al.  ·  TPAMI 2019",
             fill=WASH, fg=PURPLE, size=9)

    # Three concept cards stacked
    add_concept_card(slide, 0.36, 1.55, 9.2, 0.95,
                     num=1, header="Part Affinity Fields",
                     body="Per-pixel 2D vectors aligned along each limb's direction — the network outputs PAFs alongside keypoint heatmaps.",
                     color=PURPLE)
    add_concept_card(slide, 0.36, 2.65, 9.2, 0.95,
                     num=2, header="Bottom-up association",
                     body="Detect every keypoint first, then bipartite-match to people via PAF scores. No per-person crop.",
                     color=PURPLE_SOFT)
    add_concept_card(slide, 0.36, 3.75, 9.2, 0.95,
                     num=3, header="O(1) inference  vs  number of people",
                     body="Real-time multi-person at 22 fps — runtime independent of crowd size.",
                     color=CYAN)

    # Bottom note: limit + bridge to next paper
    add_simple_text(slide, 1.30, 5.22, 8.3, 0.24,
                    "Limit:  precision capped by heatmap resolution  →  HRNet attacks this.  HMR consumes the 2D output.",
                    font=BODY_FONT, size=10, color=GREY, italic=True)


def build_hmr(slide):
    add_title(slide, "HMR", kicker="02  ·  3D Body Recovery  ·  Single-pass regression")
    add_speaker_chip(slide, "Yiqiao")
    add_pill(slide, 7.20, 0.42, 2.40, 0.30, "Kanazawa et al.  ·  CVPR 2018",
             fill=WASH, fg=PURPLE, size=9)

    # Hero stat: 60s → 100ms
    add_simple_text(slide, 0.36, 1.55, 9.2, 0.78,
                    "60 s   →   ~100 ms",
                    font=HEAD_FONT, size=44, color=PURPLE, bold=True,
                    spacing=-40, align=PP_ALIGN.CENTER)
    add_simple_text(slide, 0.36, 2.30, 9.2, 0.30,
                    "iterative SMPLify   →   one forward pass   ·   measured on RTX PRO 6000",
                    font=BODY_FONT, size=11, color=GREY, italic=True,
                    spacing=80, align=PP_ALIGN.CENTER)

    # Three concept cards in a row
    card_w = 2.92; card_gap = 0.12; T = 2.92
    cards = [
        (1, "Single-pass regression", "Image  →  82-d SMPL params  (72 θ + 10 β) + camera.", PURPLE),
        (2, "Reprojection loss", "3D joints projected back to 2D, supervised by HRNet/OpenPose-style detections.", CYAN),
        (3, "Adversarial body prior", "Discriminator on 4M Mocap poses keeps outputs anatomically plausible.", CORAL),
    ]
    for i, (num, head, body, col) in enumerate(cards):
        L = 0.36 + i * (card_w + card_gap)
        # Subtle card background
        add_filled_rect(slide, L, T, card_w, 1.95, fill=WASH, rounded=True)
        # Numbered chip
        add_pill(slide, L + 0.20, T + 0.18, 0.40, 0.26, f"{num:02d}",
                 fill=col, fg=WHITE, size=10)
        # Header
        add_simple_text(slide, L + 0.20, T + 0.55, card_w - 0.30, 0.34,
                        head, font=HEAD_FONT, size=12, color=DARK, bold=True)
        # Body
        add_simple_text(slide, L + 0.20, T + 0.92, card_w - 0.30, 1.0,
                        body, font=BODY_FONT, size=10, color=GREY)

    add_simple_text(slide, 1.30, 5.22, 8.3, 0.24,
                    "Key insight  —  the 2D papers are not HMR's competitors;  they are its training signal.",
                    font=BODY_FONT, size=10, color=PURPLE, italic=True, spacing=80)


def build_hrnet(slide):
    add_title(slide, "HRNet", kicker="03  ·  2D Pose Estimation  ·  Top-down  ·  high resolution")
    add_speaker_chip(slide, "Taijia")
    add_pill(slide, 7.20, 0.42, 2.40, 0.30, "Sun et al.  ·  CVPR 2019",
             fill=WASH, fg=PURPLE, size=9)

    # Three concept cards, left column
    add_concept_card(slide, 0.36, 1.55, 5.0, 0.85,
                     num=1, header="Parallel multi-resolution",
                     body="Branches at 4 scales fuse continuously — high-res stream is never lost.",
                     color=PURPLE)
    add_concept_card(slide, 0.36, 2.55, 5.0, 0.85,
                     num=2, header="Top-down pipeline",
                     body="Detector first, then per-person crop  ·  cost grows with #people.",
                     color=PURPLE_SOFT)
    add_concept_card(slide, 0.36, 3.55, 5.0, 0.85,
                     num=3, header="Wins on small joints",
                     body="Crisp wrists / ankles / toes — exactly what HMR's reprojection loss leans on.",
                     color=CYAN)

    # Right-side image
    yoga = RESULTS / "2d_vis" / "img2_complex_pose.jpg"
    if yoga.exists():
        add_image(slide, yoga, left=5.55, top=1.30, height=3.4)
    add_caption(slide, 5.55, 4.78, 4.1,
                "Our HRNet output  ·  yoga photo  ·  notice the wrist / ankle precision",
                size=9, italic=True, color=GREY)

    add_simple_text(slide, 1.30, 5.22, 8.3, 0.24,
                    "Two complementary endpoints of the 2D paradigm  —  bottom-up speed  vs  top-down precision.",
                    font=BODY_FONT, size=10, color=GREY, italic=True)


def build_bridge(slide):
    add_title(slide, "From 2D to 3D", kicker="THE BRIDGE  ·  SMPL  +  Reprojection loss")
    add_speaker_chip(slide, "Taijia")

    # Top: SMPL formula in a pill-styled banner
    add_filled_rect(slide, 0.36, 1.50, 9.2, 0.62, fill=WASH, rounded=True)
    add_simple_text(slide, 0.36, 1.58, 9.2, 0.45,
                    "SMPL    M ( θ , β )    ·    72-d pose  +  10-d shape    →    6890-vertex mesh",
                    font=HEAD_FONT, size=14, color=PURPLE, bold=True,
                    spacing=80, align=PP_ALIGN.CENTER)

    # Left: vertical flow diagram (compacted to fit above the speaker chip)
    flow_x = 0.65; flow_w = 4.5
    rows = [
        ("OpenPose  +  HRNet", "produce 2D keypoints", CYAN),
        ("HMR", "predicts 3D joints", PURPLE),
        ("Project back to 2D", "weak-perspective camera", PURPLE_SOFT),
        ("L2  vs  HRNet 2D", "→  reprojection loss", CORAL),
    ]
    box_h = 0.55; gap = 0.10
    T0 = 2.32
    for i, (head, body, col) in enumerate(rows):
        T = T0 + i * (box_h + gap)
        add_filled_rect(slide, flow_x, T, flow_w, box_h, fill=WASH, rounded=True,
                        line=col, line_w=0.75)
        add_simple_text(slide, flow_x + 0.18, T + 0.04, flow_w - 0.30, 0.26,
                        head, font=HEAD_FONT, size=12, color=col, bold=True)
        add_simple_text(slide, flow_x + 0.18, T + 0.28, flow_w - 0.30, 0.24,
                        body, font=BODY_FONT, size=10, color=GREY)
        if i < len(rows) - 1:
            add_arrow(slide, flow_x + flow_w / 2, T + box_h,
                       flow_x + flow_w / 2, T + box_h + gap - 0.02,
                       color=GREY_SOFT, weight=0.75)

    # Right: empirical reprojection overlay (matches diagram height range)
    reproj = RESULTS / "extras" / "reproj_img1_standing.png"
    if reproj.exists():
        add_image(slide, reproj, left=6.95, top=2.25, height=2.45)
    add_caption(slide, 5.55, 4.74, 4.1,
                "Our HMR 3D joints (×) reprojected vs HRNet 2D (○)",
                size=9, italic=True)

    # Bottom callout — placed to the right of the speaker chip, no overlap
    add_simple_text(slide, 1.40, 5.20, 8.0, 0.24,
                    "★  The 2D papers are HMR's training signal — not its competitors.",
                    font=BODY_FONT, size=10, color=PURPLE, bold=True, spacing=80)


def build_4d(slide):
    add_title(slide, "From 3D to 4D", kicker="HMR 2.0  +  PHALP  ·  Adding the time axis")
    add_speaker_chip(slide, "Taijia")

    # Timeline with year markers
    timeline_y = 2.20
    timeline_x0, timeline_x1 = 0.80, 5.40
    add_line(slide, timeline_x0, timeline_y, timeline_x1, timeline_y,
             color=GREY_SOFT, weight=1.5)
    milestones = [
        (timeline_x0,             "2018", "HMR",       "single image"),
        ((timeline_x0 + timeline_x1) / 2, "2022", "PHALP",     "video tracking"),
        (timeline_x1,             "2023", "HMR 2.0",   "ViT backbone"),
    ]
    for x, year, name, sub in milestones:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(x - 0.07), Inches(timeline_y - 0.07),
                                     Inches(0.14), Inches(0.14))
        dot.fill.solid(); dot.fill.fore_color.rgb = PURPLE
        dot.line.color.rgb = PURPLE
        add_simple_text(slide, x - 0.55, timeline_y - 0.55, 1.10, 0.24,
                        year, font=BODY_FONT, size=10, color=GREY,
                        bold=True, spacing=200, align=PP_ALIGN.CENTER)
        add_simple_text(slide, x - 0.95, timeline_y + 0.16, 1.90, 0.34,
                        name, font=HEAD_FONT, size=15, color=PURPLE,
                        bold=True, align=PP_ALIGN.CENTER)
        add_simple_text(slide, x - 1.05, timeline_y + 0.52, 2.10, 0.26,
                        sub, font=BODY_FONT, size=9, color=GREY,
                        italic=True, align=PP_ALIGN.CENTER)

    # Two takeaways below the timeline
    add_filled_rect(slide, 0.36, 3.60, 5.50, 1.30, fill=WASH, rounded=True)
    add_simple_text(slide, 0.55, 3.70, 5.20, 0.34,
                    "4D  =  3D  +  time", font=HEAD_FONT, size=18,
                    color=PURPLE, bold=True, spacing=-15)
    add_simple_text(slide, 0.55, 4.05, 5.20, 0.85,
                    "No new mechanism — per-frame HMR 2.0 stitched by PHALP's tracking. "
                    "Output isn't a video; it's a per-frame SMPL time-series ready for biomechanics.",
                    font=BODY_FONT, size=10, color=DARK)

    # Right: tennis frame as evidence
    frame = RESULTS / "tennis_phalp" / "frame_2.0s.png"
    if frame.exists():
        add_image(slide, frame, left=6.10, top=1.75, height=3.15)
    add_caption(slide, 6.10, 4.95, 3.50,
                "PHALP frame @ 2.0 s  ·  HMR 2.0 mesh on track 1", size=9, italic=True)


def build_image_demo(slide):
    add_title(slide, "Image pipeline", kicker="LIVE DEMO  ·  4 in-the-wild photos")
    add_speaker_chip(slide, "Taijia")

    # Subtitle banner
    add_simple_text(slide, 0.36, 1.18, 9.2, 0.26,
                    "each panel:  input  →  HRNet 2D  →  HMR 2.0 (front)  →  HMR 2.0 (side)   ·   ~100 ms / image",
                    font=BODY_FONT, size=10, color=GREY,
                    italic=True, spacing=80, align=PP_ALIGN.CENTER)

    # 4 images in a row, each ≈ 2.2" square (quadplot is ≈ 1:1)
    quads = [
        ("quadplot_img1_standing.png",     "Standing",      "baseline",         PURPLE),
        ("quadplot_img2_complex_pose.png", "Yoga",          "prior holds",      CORAL),
        ("quadplot_img3_occluded.png",     "Occluded",      "prior fills",      CYAN),
        ("quadplot_img4_multi_person.png", "Multi-person",  "top-down crop",    PURPLE_SOFT),
    ]
    L0 = 0.42; img_w = 2.20; gap = 0.10; img_top = 1.65
    for i, (fname, head, sub, col) in enumerate(quads):
        path = RESULTS / fname
        if not path.exists():
            continue
        L = L0 + i * (img_w + gap)
        add_image(slide, path, left=L, top=img_top, width=img_w)
        # Pill caption below
        add_pill(slide, L + 0.12, img_top + img_w + 0.10, img_w - 0.24, 0.28,
                 head, fill=col, fg=WHITE, size=10)
        add_simple_text(slide, L, img_top + img_w + 0.42, img_w, 0.24,
                        sub, font=BODY_FONT, size=9, color=GREY,
                        italic=True, align=PP_ALIGN.CENTER)

    add_simple_text(slide, 1.30, 5.22, 8.3, 0.24,
                    "Same forward pass.  No per-image tuning.  No iteration.",
                    font=BODY_FONT, size=10, color=PURPLE, bold=True,
                    spacing=200, align=PP_ALIGN.CENTER)


def build_video_demo(slide):
    add_title(slide, "Video pipeline", kicker="LIVE DEMO  ·  8.6 s tennis serve")
    add_speaker_chip(slide, "Taijia")

    # Try to embed video; fall back to static frame
    video = RESULTS / "extras" / "tennis_clip.mp4"
    poster = RESULTS / "tennis_phalp" / "frame_2.0s.png"
    embedded = False
    if video.exists() and poster.exists():
        try:
            slide.shapes.add_movie(str(video), Inches(0.36), Inches(1.50),
                                   Inches(5.4), Inches(3.05),
                                   poster_frame_image=str(poster))
            embedded = True
        except Exception:
            embedded = False
    if not embedded and poster.exists():
        add_image(slide, poster, left=0.36, top=1.50, width=5.4, height=3.05)
    add_caption(slide, 0.36, 4.62, 5.4,
                "tennis_clip.mp4  ·  5 s loop  ·  double-click in PowerPoint to play",
                size=9, italic=True)

    # Right column: hero stats + numbered notes
    add_simple_text(slide, 6.05, 1.50, 3.55, 0.46,
                    "216 frames", font=HEAD_FONT, size=24, color=PURPLE,
                    bold=True, spacing=-15)
    add_simple_text(slide, 6.05, 1.95, 3.55, 0.24,
                    "every frame  ·  one HMR 2.0 forward pass",
                    font=BODY_FONT, size=10, color=GREY, italic=True)

    # Three small numbered points
    pts = [
        ("Same HMR 2.0 used on each frame", PURPLE),
        ("Tracking re-IDs the player across motion blur", CYAN),
        ("Output: structured 3D time-series — biomechanics-ready", CORAL),
    ]
    for i, (text, col) in enumerate(pts):
        y = 2.50 + i * 0.62
        add_filled_rect(slide, 6.05, y, 0.24, 0.24, fill=col, rounded=True)
        add_simple_text(slide, 6.05, y + 0.02, 0.24, 0.22,
                        f"{i+1}", font=HEAD_FONT, size=11, color=WHITE,
                        bold=True, align=PP_ALIGN.CENTER)
        add_simple_text(slide, 6.36, y, 3.30, 0.55,
                        text, font=BODY_FONT, size=10, color=DARK)

    add_simple_text(slide, 1.30, 5.22, 8.3, 0.24,
                    "The valuable output isn't the video  —  it's the data behind it.",
                    font=BODY_FONT, size=10, color=PURPLE, bold=True,
                    spacing=200, align=PP_ALIGN.CENTER)


def build_analysis(slide):
    """Fills the existing 'Analysis, Limitations & Next Steps' section divider.

    Replaces three columns of bullets with three horizontal cards, each
    headed with a colored vertical accent bar + label.
    """
    add_speaker_chip(slide, "Taijia")
    add_simple_text(slide, 0.36, 1.18, 9.2, 0.30,
                    "Three take-aways  ·  three honest limitations  ·  three forward directions",
                    font=BODY_FONT, size=11, color=GREY,
                    italic=True, spacing=80, align=PP_ALIGN.CENTER)

    rows = [
        ("01", "TAKE-AWAYS", PURPLE, [
            "2D papers supervise the 3D paper.",
            "Single-pass paradigm wins everywhere except absolute precision.",
            "Adding the time axis is the natural next step, not a new field.",
        ]),
        ("02", "LIMITATIONS  TODAY", CORAL, [
            "Per-frame jitter and global scale drift in video.",
            "3D ground truth is scarce — in-the-wild domains are narrow.",
            "Mesh surfaces lack clothing, hair, soft tissue.",
        ]),
        ("03", "WHAT COMES NEXT", CYAN, [
            "VIBE / SLAHMR  —  temporal smoothing + global camera.",
            "Sapiens  (Meta 2024)  —  human foundation model.",
            "3DGS-Avatar  —  neural / Gaussian surfaces.",
        ]),
    ]
    T0 = 1.62; row_h = 1.05; row_gap = 0.10
    for i, (num, label, col, items) in enumerate(rows):
        T = T0 + i * (row_h + row_gap)
        # Vertical accent bar
        add_filled_rect(slide, 0.36, T, 0.10, row_h, fill=col)
        # Number + label
        add_simple_text(slide, 0.55, T + 0.04, 0.50, 0.30,
                        num, font=HEAD_FONT, size=14, color=col, bold=True)
        add_simple_text(slide, 1.05, T + 0.06, 4.0, 0.28,
                        label, font=BODY_FONT, size=11, color=col,
                        bold=True, spacing=200)
        # Three inline items
        for j, item in enumerate(items):
            y = T + 0.40 + j * 0.22
            add_simple_text(slide, 0.55, y, 9.0, 0.22,
                            "•   " + item, font=BODY_FONT, size=10, color=DARK)


# ────────────────────────── orchestration ──────────────────────────

def fix_title_placeholder(slide, new_title="From 2D Keypoints to 4D Humans"):
    """Replace the template's placeholder ('112233') with the real project title.

    Also restyles the title using our new font hierarchy.
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip() == "112233":
                    r.text = new_title
                    set_run(r, font=HEAD_FONT, size=36, color=PURPLE,
                            bold=True, spacing=-25)
                    return True
    return False


def restyle_section_dividers(*slides):
    """Apply the new font to existing template section-divider titles."""
    for slide in slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip() and any(c.isalpha() for c in r.text):
                        # rewrite the run with our font
                        set_run(r, font=HEAD_FONT, size=28, color=PURPLE,
                                bold=True, spacing=-20)


def reset_to_template(prs):
    """Strip any previously-generated content so the deck is idempotent.

    Fully removes BLANK-layout slides (sldIdLst entry + rel + package part)
    and drops every shape after the first (title) on each TITLE_AND_BODY
    divider. Without removing the package part you get duplicate XML names
    on save, which corrupts the .pptx.
    """
    sldIdLst = prs.slides._sldIdLst
    rels = prs.part.rels
    package = prs.part.package

    # Find and fully drop BLANK slides
    to_remove = []
    for sldId in list(sldIdLst):
        rId = sldId.rId
        slide_part = rels[rId].target_part
        if slide_part.slide.slide_layout.name == "BLANK":
            to_remove.append((sldId, rId, slide_part))

    for sldId, rId, slide_part in to_remove:
        sldIdLst.remove(sldId)
        prs.part.drop_rel(rId)
        # Remove from package's internal parts dict + iter_parts cache
        pkg_parts = package.__dict__.get("_parts", None)
        if pkg_parts is None:
            # python-pptx stores parts via the iter_parts() generator backed
            # by the relationships graph. Once both the rel from the
            # presentation part and the slide rels themselves are dropped,
            # the part becomes unreachable and won't be serialized.
            pass
        # Drop slide's own relationships (e.g. slideLayout) so it's fully
        # disconnected from the relationship graph.
        for slide_rId in list(slide_part.rels):
            slide_part.drop_rel(slide_rId)

    if to_remove:
        print(f"  reset: removed {len(to_remove)} previously-added BLANK slide(s)")

    # Strip extra shapes from section-divider slides
    for slide in prs.slides:
        if slide.slide_layout.name == "TITLE_AND_BODY" and len(slide.shapes) > 1:
            for shape in list(slide.shapes)[1:]:
                slide.shapes._spTree.remove(shape._element)
            print(f"  reset: cleared extra shapes from '{slide.shapes[0].text_frame.text[:30]}…'")


def main():
    if not TEMPLATE.exists():
        raise RuntimeError(
            f"{TEMPLATE} missing — run presentation/_make_template.py once to "
            "extract a clean 5-slide template from the current Intro to CV.pptx."
        )
    prs = Presentation(TEMPLATE)
    layout_blank = prs.slide_layouts[15]

    # Strip any stale per-slide content carried over from previous builds.
    # (The template.pptx may have legacy shapes on the Analysis divider.)
    for slide in prs.slides:
        if slide.slide_layout.name == "TITLE_AND_BODY" and len(slide.shapes) > 1:
            for shape in list(slide.shapes)[1:]:
                slide.shapes._spTree.remove(shape._element)

    sldIdLst = prs.slides._sldIdLst
    orig_sldIds = list(sldIdLst)
    if len(orig_sldIds) != 5:
        raise RuntimeError(f"expected 5 template slides, found {len(orig_sldIds)}")

    s_title, s_rw, s_pr, s_an, s_thanks = list(prs.slides)
    print(f"Loaded template: {len(orig_sldIds)} slides")

    if fix_title_placeholder(s_title):
        print("✓ replaced title-slide placeholder + restyled title")

    # Restyle the 3 section-divider titles to match the new typography
    restyle_section_dividers(s_rw, s_pr, s_an)
    print("✓ restyled 3 section-divider titles")

    # Add the 7 new content slides (appended).
    builders = [build_openpose, build_hmr, build_hrnet, build_bridge,
                build_4d, build_image_demo, build_video_demo]
    for builder in builders:
        slide = prs.slides.add_slide(layout_blank)
        builder(slide)

    new_sldIds = list(sldIdLst)[5:]
    assert len(new_sldIds) == 7

    # Fill the Analysis section divider's body
    build_analysis(s_an)

    # Reorder: title, RW, [op, hmr, hrnet, bridge, 4d], PR, [imdemo, viddemo], AN, Thanks
    desired_xml = [
        orig_sldIds[0],                       # title
        orig_sldIds[1],                       # RW divider
        new_sldIds[0], new_sldIds[1], new_sldIds[2], new_sldIds[3], new_sldIds[4],
        orig_sldIds[2],                       # PR divider
        new_sldIds[5], new_sldIds[6],
        orig_sldIds[3], orig_sldIds[4],       # Analysis, Thanks
    ]
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for el in desired_xml:
        sldIdLst.append(el)

    prs.save(OUT)
    dedupe_zip(OUT)
    print(f"✓ wrote {OUT.relative_to(ROOT)}  ·  {len(prs.slides)} slides total")


def dedupe_zip(path: Path):
    """Rewrite the .pptx so only slides referenced by presentation.xml are kept.

    python-pptx's slide-deletion leaves orphan slide parts in the package; on
    save these collide with new slides on the same filename. The collision
    breaks LibreOffice and confuses PowerPoint's slide ordering. We resolve
    it by walking the official slide list (from presentation.xml + its rels)
    and keeping only those slide files in the output zip.
    """
    import shutil
    import zipfile
    from xml.etree import ElementTree as ET

    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"

    with zipfile.ZipFile(path, "r") as zin:
        names = zin.namelist()

        # 1. Read the slide order from presentation.xml + its rels.
        pres_xml = ET.fromstring(zin.read("ppt/presentation.xml"))
        sldIdLst = pres_xml.find(f"{{{P_NS}}}sldIdLst")
        rIds_in_order = [s.get(f"{{{R_NS}}}id") for s in sldIdLst.findall(f"{{{P_NS}}}sldId")]

        rels_xml = ET.fromstring(zin.read("ppt/_rels/presentation.xml.rels"))
        rId_to_target = {}
        for rel in rels_xml.findall(f"{{{REL_NS}}}Relationship"):
            rId_to_target[rel.get("Id")] = rel.get("Target")

        # The official ordered list of slide partnames the deck actually uses.
        live_slides = []
        for rId in rIds_in_order:
            target = rId_to_target.get(rId, "")
            if target:
                # targets are relative to ppt/, so e.g. "slides/slide3.xml"
                live_slides.append("ppt/" + target)
        live_slide_set = set(live_slides)

        # Anything under ppt/slides/ that isn't in live_slides is an orphan.
        orphans = set()
        for n in names:
            if n.startswith("ppt/slides/slide") and n.endswith(".xml") and n not in live_slide_set:
                orphans.add(n)
                # Also drop its rels file
                rel_name = n.replace("slides/", "slides/_rels/") + ".rels"
                orphans.add(rel_name)

        # 2. Rewrite the zip — drop orphans, dedupe duplicate names by last-write.
        last_index = {}
        for i, info in enumerate(zin.infolist()):
            if info.filename in orphans:
                continue
            last_index[info.filename] = i

        tmp_path = path.with_suffix(".tmp.pptx")
        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            written = set()
            for i, info in enumerate(zin.infolist()):
                if info.filename in orphans:
                    continue
                if last_index.get(info.filename) != i:
                    continue
                if info.filename in written:
                    continue
                written.add(info.filename)
                zout.writestr(info, zin.read(info.filename))

    shutil.move(tmp_path, path)


if __name__ == "__main__":
    main()
